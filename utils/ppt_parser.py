"""
Utility to parse PPT/PPTX files and extract text from slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os


def extract_text_from_pptx(file_path):
    """
    Extract all text content from a PPTX file.

    Args:
        file_path (str): Path to the PPTX file.

    Returns:
        dict: Contains 'slides' (list of slide texts), 'slide_count' (int),
              'file_name' (str), 'file_size' (str)
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    if not file_path.lower().endswith(('.ppt', '.pptx')):
        raise ValueError("File must be a PPT or PPTX file.")

    prs = Presentation(file_path)
    slides = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        slide_text.append(" | ".join(row_text))
        slides.append({
            "slide_number": slide_num,
            "text": "\n".join(slide_text) if slide_text else "[No extractable text]"
        })

    file_size = os.path.getsize(file_path)
    file_size_str = _format_file_size(file_size)
    file_name = os.path.basename(file_path)

    return {
        "slides": slides,
        "slide_count": len(slides),
        "file_name": file_name,
        "file_size": file_size_str
    }


def get_text_preview(slides, max_chars=500):
    """
    Get a preview of extracted text from slides.

    Args:
        slides (list): List of slide dicts from extract_text_from_pptx.
        max_chars (int): Maximum characters for preview.

    Returns:
        str: Preview text.
    """
    all_text = []
    total_chars = 0
    for slide in slides:
        text = f"--- Slide {slide['slide_number']} ---\n{slide['text']}"
        all_text.append(text)
        total_chars += len(text)
        if total_chars > max_chars:
            remaining = max_chars - sum(len(t) for t in all_text[:-1])
            all_text[-1] = text[:remaining] + "\n... [truncated]"
            break

    return "\n\n".join(all_text)


def get_combined_text(slides):
    """
    Combine all slide texts into a single string for quiz generation.

    Args:
        slides (list): List of slide dicts from extract_text_from_pptx.

    Returns:
        str: Combined text.
    """
    parts = []
    for slide in slides:
        if slide["text"] != "[No extractable text]":
            parts.append(slide["text"])
    return "\n\n".join(parts)


def _format_file_size(size_bytes):
    """Format file size in human-readable format."""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes / (1024 * 1024):.1f} MB"
    else:
        return f"{size_bytes / (1024 * 1024 * 1024):.1f} GB"